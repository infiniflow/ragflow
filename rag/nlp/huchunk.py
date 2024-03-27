#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
import re
import os
import copy
import base64
import magic
from dataclasses import dataclass
from typing import List
import numpy as np
from io import BytesIO


class HuChunker:

    @dataclass
    class Fields:
        text_chunks: List = None
        table_chunks: List = None

    def __init__(self):
        self.MAX_LVL = 12
        self.proj_patt = [
            (r"第[零一二三四五六七八九十百]+章", 1),
            (r"第[零一二三四五六七八九十百]+[条节]", 2),
            (r"[零一二三四五六七八九十百]+[、 　]", 3),
            (r"[\(（][零一二三四五六七八九十百]+[）\)]", 4),
            (r"[0-9]+(、|\.[　 ]|\.[^0-9])", 5),
            (r"[0-9]+\.[0-9]+(、|[ 　]|[^0-9])", 6),
            (r"[0-9]+\.[0-9]+\.[0-9]+(、|[ 　]|[^0-9])", 7),
            (r"[0-9]+\.[0-9]+\.[0-9]+\.[0-9]+(、|[ 　]|[^0-9])", 8),
            (r".{,48}[：:?？]@", 9),
            (r"[0-9]+）", 10),
            (r"[\(（][0-9]+[）\)]", 11),
            (r"[零一二三四五六七八九十百]+是", 12),
            (r"[⚫•➢✓ ]", 12)
        ]
        self.lines = []

    def _garbage(self, txt):
        patt = [
            r"(在此保证|不得以任何形式翻版|请勿传阅|仅供内部使用|未经事先书面授权)",
            r"(版权(归本公司)*所有|免责声明|保留一切权力|承担全部责任|特别声明|报告中涉及)",
            r"(不承担任何责任|投资者的通知事项：|任何机构和个人|本报告仅为|不构成投资)",
            r"(不构成对任何个人或机构投资建议|联系其所在国家|本报告由从事证券交易)",
            r"(本研究报告由|「认可投资者」|所有研究报告均以|请发邮件至)",
            r"(本报告仅供|市场有风险，投资需谨慎|本报告中提及的)",
            r"(本报告反映|此信息仅供|证券分析师承诺|具备证券投资咨询业务资格)",
            r"^(时间|签字|签章)[:：]",
            r"(参考文献|目录索引|图表索引)",
            r"[ ]*年[ ]+月[ ]+日",
            r"^(中国证券业协会|[0-9]+年[0-9]+月[0-9]+日)$",
            r"\.{10,}",
            r"(———————END|帮我转发|欢迎收藏|快来关注我吧)"
        ]
        return any([re.search(p, txt) for p in patt])

    def _proj_match(self, line):
        for p, j in self.proj_patt:
            if re.match(p, line):
                return j
        return

    def _does_proj_match(self):
        mat = [None for _ in range(len(self.lines))]
        for i in range(len(self.lines)):
            mat[i] = self._proj_match(self.lines[i])
        return mat

    def naive_text_chunk(self, text, ti="", MAX_LEN=612):
        if text:
            self.lines = [l.strip().replace(u'\u3000', u' ')
                          .replace(u'\xa0', u'')
                          for l in text.split("\n\n")]
            self.lines = [l for l in self.lines if not self._garbage(l)]
            self.lines = [re.sub(r"([ ]+|&nbsp;)", " ", l)
                          for l in self.lines if l]
        if not self.lines:
            return []
        arr = self.lines

        res = [""]
        i = 0
        while i < len(arr):
            a = arr[i]
            if not a:
                i += 1
                continue
            if len(a) > MAX_LEN:
                a_ = a.split("\n")
                if len(a_) >= 2:
                    arr.pop(i)
                    for j in range(2, len(a_) + 1):
                        if len("\n".join(a_[:j])) >= MAX_LEN:
                            arr.insert(i, "\n".join(a_[:j - 1]))
                            arr.insert(i + 1, "\n".join(a_[j - 1:]))
                            break
                    else:
                        assert False, f"Can't split: {a}"
                    continue

            if len(res[-1]) < MAX_LEN / 3:
                res[-1] += "\n" + a
            else:
                res.append(a)
            i += 1

        if ti:
            for i in range(len(res)):
                if res[i].find("——来自") >= 0:
                    continue
                res[i] += f"\t——来自“{ti}”"

        return res

    def _merge(self):
        # merge continuous same level text
        lines = [self.lines[0]] if self.lines else []
        for i in range(1, len(self.lines)):
            if self.mat[i] == self.mat[i - 1] \
               and len(lines[-1]) < 256 \
               and len(self.lines[i]) < 256:
                lines[-1] += "\n" + self.lines[i]
                continue
            lines.append(self.lines[i])
        self.lines = lines
        self.mat = self._does_proj_match()
        return self.mat

    def text_chunks(self, text):
        if text:
            self.lines = [l.strip().replace(u'\u3000', u' ')
                          .replace(u'\xa0', u'')
                          for l in re.split(r"[\r\n]", text)]
            self.lines = [l for l in self.lines if not self._garbage(l)]
            self.lines = [l for l in self.lines if l]
        self.mat = self._does_proj_match()
        mat = self._merge()

        tree = []
        for i in range(len(self.lines)):
            tree.append({"proj": mat[i],
                         "children": [],
                         "read": False})
        # find all children
        for i in range(len(self.lines) - 1):
            if tree[i]["proj"] is None:
                continue
            ed = i + 1
            while ed < len(tree) and (tree[ed]["proj"] is None or
                                      tree[ed]["proj"] > tree[i]["proj"]):
                ed += 1

            nxt = tree[i]["proj"] + 1
            st = set([p["proj"] for p in tree[i + 1: ed] if p["proj"]])
            while nxt not in st:
                nxt += 1
                if nxt > self.MAX_LVL:
                    break
            if nxt <= self.MAX_LVL:
                for j in range(i + 1, ed):
                    if tree[j]["proj"] is not None:
                        break
                    tree[i]["children"].append(j)
                for j in range(i + 1, ed):
                    if tree[j]["proj"] != nxt:
                        continue
                    tree[i]["children"].append(j)
            else:
                for j in range(i + 1, ed):
                    tree[i]["children"].append(j)

        # get DFS combinations, find all the paths to leaf
        paths = []

        def dfs(i, path):
            nonlocal tree, paths
            path.append(i)
            tree[i]["read"] = True
            if len(self.lines[i]) > 256:
                paths.append(path)
                return
            if not tree[i]["children"]:
                if len(path) > 1 or len(self.lines[i]) >= 32:
                    paths.append(path)
                return
            for j in tree[i]["children"]:
                dfs(j, copy.deepcopy(path))

        for i, t in enumerate(tree):
            if t["read"]:
                continue
            dfs(i, [])

        # concat txt on the path for all paths
        res = []
        lines = np.array(self.lines)
        for p in paths:
            if len(p) < 2:
                tree[p[0]]["read"] = False
                continue
            txt = "\n".join(lines[p[:-1]]) + "\n" + lines[p[-1]]
            res.append(txt)
        # concat continuous orphans
        assert len(tree) == len(lines)
        ii = 0
        while ii < len(tree):
            if tree[ii]["read"]:
                ii += 1
                continue
            txt = lines[ii]
            e = ii + 1
            while e < len(tree) and not tree[e]["read"] and len(txt) < 256:
                txt += "\n" + lines[e]
                e += 1
            res.append(txt)
            ii = e

        # if the node has not been read, find its daddy
        def find_daddy(st):
            nonlocal lines, tree
            proj = tree[st]["proj"]
            if len(self.lines[st]) > 512:
                return [st]
            if proj is None:
                proj = self.MAX_LVL + 1
            for i in range(st - 1, -1, -1):
                if tree[i]["proj"] and tree[i]["proj"] < proj:
                    a = [st] + find_daddy(i)
                    return a
            return []

        return res


class PdfChunker(HuChunker):

    def __init__(self, pdf_parser):
        self.pdf = pdf_parser
        super().__init__()

    def tableHtmls(self, pdfnm):
        _, tbls = self.pdf(pdfnm, return_html=True)
        res = []
        for img, arr in tbls:
            if arr[0].find("<table>") < 0:
                continue
            buffered = BytesIO()
            if img:
                img.save(buffered, format="JPEG")
            img_str = base64.b64encode(
                buffered.getvalue()).decode('utf-8') if img else ""
            res.append({"table": arr[0], "image": img_str})
        return res

    def html(self, pdfnm):
        txts, tbls = self.pdf(pdfnm, return_html=True)
        res = []
        txt_cks = self.text_chunks(txts)
        for txt, img in [(self.pdf.remove_tag(c), self.pdf.crop(c))
                         for c in txt_cks]:
            buffered = BytesIO()
            if img:
                img.save(buffered, format="JPEG")
            img_str = base64.b64encode(
                buffered.getvalue()).decode('utf-8') if img else ""
            res.append({"table": "<p>%s</p>" % txt.replace("\n", "<br/>"),
                        "image": img_str})

        for img, arr in tbls:
            if not arr:
                continue
            buffered = BytesIO()
            if img:
                img.save(buffered, format="JPEG")
            img_str = base64.b64encode(
                buffered.getvalue()).decode('utf-8') if img else ""
            res.append({"table": arr[0], "image": img_str})

        return res

    def __call__(self, pdfnm, return_image=True, naive_chunk=False):
        flds = self.Fields()
        text, tbls = self.pdf(pdfnm)
        fnm = pdfnm
        txt_cks = self.text_chunks(text) if not naive_chunk else \
            self.naive_text_chunk(text, ti=fnm if isinstance(fnm, str) else "")
        flds.text_chunks = [(self.pdf.remove_tag(c),
                             self.pdf.crop(c) if return_image else None) for c in txt_cks]

        flds.table_chunks = [(arr, img if return_image else None)
                             for img, arr in tbls]
        return flds


class DocxChunker(HuChunker):

    def __init__(self, doc_parser):
        self.doc = doc_parser
        super().__init__()

    def _does_proj_match(self):
        mat = []
        for s in self.styles:
            s = s.split(" ")[-1]
            try:
                mat.append(int(s))
            except Exception as e:
                mat.append(None)
        return mat

    def _merge(self):
        i = 1
        while i < len(self.lines):
            if self.mat[i] == self.mat[i - 1] \
               and len(self.lines[i - 1]) < 256 \
               and len(self.lines[i]) < 256:
                self.lines[i - 1] += "\n" + self.lines[i]
                self.styles.pop(i)
                self.lines.pop(i)
                self.mat.pop(i)
                continue
            i += 1
        self.mat = self._does_proj_match()
        return self.mat

    def __call__(self, fnm):
        flds = self.Fields()
        flds.title = os.path.splitext(
            os.path.basename(fnm))[0] if isinstance(
            fnm, type("")) else ""
        secs, tbls = self.doc(fnm)
        self.lines = [l for l, s in secs]
        self.styles = [s for l, s in secs]

        txt_cks = self.text_chunks("")
        flds.text_chunks = [(t, None) for t in txt_cks if not self._garbage(t)]
        flds.table_chunks = [(tb, None) for tb in tbls for t in tb if t]
        return flds


class ExcelChunker(HuChunker):

    def __init__(self, excel_parser):
        self.excel = excel_parser
        super().__init__()

    def __call__(self, fnm):
        flds = self.Fields()
        flds.text_chunks = [(t, None) for t in self.excel(fnm)]
        flds.table_chunks = []
        return flds


class PptChunker(HuChunker):

    def __init__(self):
        super().__init__()

    def __extract(self, shape):
        if shape.shape_type == 19:
            tb = shape.table
            rows = []
            for i in range(1, len(tb.rows)):
                rows.append("; ".join([tb.cell(
                    0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
            return "\n".join(rows)

        if shape.has_text_frame:
            return shape.text_frame.text

        if shape.shape_type == 6:
            texts = []
            for p in shape.shapes:
                t = self.__extract(p)
                if t:
                    texts.append(t)
            return "\n".join(texts)

    def __call__(self, fnm):
        from pptx import Presentation
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        for slide in ppt.slides:
            texts = []
            for shape in slide.shapes:
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        import aspose.slides as slides
        import aspose.pydrawing as drawing
        imgs = []
        with slides.Presentation(BytesIO(fnm)) as presentation:
            for slide in presentation.slides:
                buffered = BytesIO()
                slide.get_thumbnail(
                    0.5, 0.5).save(
                    buffered, drawing.imaging.ImageFormat.jpeg)
                imgs.append(buffered.getvalue())
        assert len(imgs) == len(
            txts), "Slides text and image do not match: {} vs. {}".format(len(imgs), len(txts))

        flds = self.Fields()
        flds.text_chunks = [(txts[i], imgs[i]) for i in range(len(txts))]
        flds.table_chunks = []

        return flds


class TextChunker(HuChunker):

    @dataclass
    class Fields:
        text_chunks: List = None
        table_chunks: List = None

    def __init__(self):
        super().__init__()

    @staticmethod
    def is_binary_file(file_path):
        mime = magic.Magic(mime=True)
        if isinstance(file_path, str):
            file_type = mime.from_file(file_path)
        else:
            file_type = mime.from_buffer(file_path)
        if 'text' in file_type:
            return False
        else:
            return True

    def __call__(self, fnm):
        flds = self.Fields()
        if self.is_binary_file(fnm):
            return flds
        txt = ""
        if isinstance(fnm, str):
            with open(fnm, "r") as f:
                txt = f.read()
        else:
            txt = fnm.decode("utf-8")
        flds.text_chunks = [(c, None) for c in self.naive_text_chunk(txt)]
        flds.table_chunks = []
        return flds


if __name__ == "__main__":
    import sys
    sys.path.append(os.path.dirname(__file__) + "/../")
    if sys.argv[1].split(".")[-1].lower() == "pdf":
        from deepdoc.parser import PdfParser
        ckr = PdfChunker(PdfParser())
    if sys.argv[1].split(".")[-1].lower().find("doc") >= 0:
        from deepdoc.parser import DocxParser
        ckr = DocxChunker(DocxParser())
    if sys.argv[1].split(".")[-1].lower().find("xlsx") >= 0:
        from deepdoc.parser import ExcelParser
        ckr = ExcelChunker(ExcelParser())

    # ckr.html(sys.argv[1])
    print(ckr(sys.argv[1]))
