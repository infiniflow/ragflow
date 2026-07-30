#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将 DeepDoc 解析链路和定制点内容添加到现有 PPT 中
从第16页开始添加
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def main():
    # 打开现有PPT
    prs = Presentation('E:/桌面/work/Presentation/DeepDoc_Architecture.pptx')

    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"Slide size: {slide_width} x {slide_height}")

    # 颜色定义
    DARK_BLUE = RGBColor(0, 51, 102)
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    GRAY = RGBColor(100, 100, 100)
    ORANGE = RGBColor(230, 126, 34)
    GREEN = RGBColor(39, 174, 96)
    PURPLE = RGBColor(142, 68, 173)
    RED = RGBColor(192, 57, 43)
    LIGHT_BG = RGBColor(248, 249, 250)

    def add_title_slide(prs, title_text, subtitle_text=""):
        """添加标题页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # 背景色块
        shape = slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        if subtitle_text:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.alignment = PP_ALIGN.CENTER

        return slide

    def add_content_slide(prs, title, content_lines):
        """添加内容页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题栏背景
        title_shape = slide.shapes.add_shape(1, 0, 0, slide_width, Inches(1.2))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.fill.background()

        # 标题文字
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 内容区域
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 处理缩进和样式
            stripped = line.lstrip()
            indent_level = (len(line) - len(stripped)) // 4

            p.text = stripped
            p.font.size = Pt(16)
            p.font.color.rgb = BLACK

            # 根据内容类型设置样式
            if stripped.startswith(('1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.')):
                p.font.bold = True
                p.font.size = Pt(18)
                p.font.color.rgb = DARK_BLUE
            elif stripped.startswith('**') and stripped.endswith('**'):
                p.font.bold = True
                p.font.size = Pt(18)
                p.font.color.rgb = ORANGE
            elif 'RAGFlowPdfParser' in stripped or 'MinerUParser' in stripped or 'OpenDataLoaderParser' in stripped:
                p.font.bold = True
                p.font.color.rgb = PURPLE
            elif 'def ' in stripped or 'class ' in stripped:
                p.font.name = 'Courier New'
                p.font.color.rgb = RGBColor(0, 100, 0)

            # 设置缩进
            if indent_level > 0:
                p.level = min(indent_level, 3)

        return slide

    def add_table_slide(prs, title, headers, rows):
        """添加表格页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 标题栏
        title_shape = slide.shapes.add_shape(1, 0, 0, slide_width, Inches(1.2))
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = DARK_BLUE
        title_shape.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 添加表格
        num_rows = len(rows) + 1
        num_cols = len(headers)
        table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6 * num_rows)).table

        # 设置表头
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE

        # 设置数据行
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(cell_text)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG

        return slide

    # ==================== 开始添加新页面 ====================
    print("开始添加新页面...")

    # 第16页：章节标题 - 解析链路
    add_title_slide(prs, "DeepDoc 解析链路", "从文件上传到结构化数据的完整流程")
    print("添加第16页：章节标题 - 解析链路")

    # 第17页：解析链路全景图
    slide17 = add_content_slide(prs, "解析链路全景图", [
        "**第一层：入口调度（rag/app/naive.py）**",
        "    - chunk() 函数：主入口，根据文件扩展名选择解析器",
        "    - PDF 文件根据 parser_config[layout_recognize] 选择解析器",
        "    - 调用对应的 by_* 函数进行解析",
        "",
        "**第二层：解析器调度层（PARSERS 字典）**",
        "    - PARSERS = {deepdoc, mineru, docling, opendataloader, ...}",
        "    - 通过名称映射到具体的 by_* 函数",
        "",
        "**第三层：具体解析器实现**",
        "    - RAGFlowPdfParser, MinerUParser, OpenDataLoaderParser 等",
        "    - 每个解析器实现 parse_pdf() 方法，返回标准格式",
        "",
        "**第四层：分块处理（Chunking）**",
        "    - naive_merge() / naive_merge_with_images()",
        "    - tokenize_chunks() 生成最终的 chunks"
    ])
    print("添加第17页：解析链路全景图")

    # 第18页：解析器对比表格
    add_table_slide(prs, "解析器能力对比",
        ["能力", "RAGFlowPdfParser", "MinerUParser", "OpenDataLoaderParser"],
        [
            ["运行方式", "本地", "远程服务", "远程服务"],
            ["OCR", "本地 PaddleOCR/ONNX", "MinerU 服务", "Docling 服务"],
            ["布局识别", "本地 YOLO/ONNX", "MinerU 服务", "Docling 服务"],
            ["表格识别", "本地 TSR 模型", "MinerU 服务", "Docling 服务"],
            ["段落合并", "XGBoost 模型", "MinerU 服务", "Docling 服务"],
            ["多后端支持", "否", "是 (pipeline/VLM)", "否"],
            ["跨页表格", "有限", "是 (middle.json)", "否"],
            ["VLM增强", "是 (VisionParser)", "是", "否"],
            ["乱码处理", "多策略", "由服务处理", "由服务处理"],
            ["中文优化", "深度", "一般", "一般"],
        ]
    )
    print("添加第18页：解析器对比表格")

    # 第19页：RAGFlowPdfParser 解析流程
    slide19 = add_content_slide(prs, "RAGFlowPdfParser 解析流程（本地）", [
        "**1. __images__()** - PDF 转图片 + 字符提取 + 乱码检测",
        "**2. _layouts_rec()** - 布局识别（文本/表格/图片区域）",
        "**3. _table_transformer_job()** - 表格结构识别 + 自动旋转",
        "**4. _text_merge()** - 文本合并（KMeans列检测 + 水平合并）",
        "**5. _concat_downward()** - 段落连接（XGBoost模型判断）",
        "**6. _filter_forpages()** - 过滤目录页等",
        "**7. _extract_table_figure()** - 提取表格和图片",
        "",
        "**核心技术栈：**",
        "    - ONNX Runtime (CPU/GPU) 或 Ascend NPU",
        "    - PaddleOCR / ONNX OCR 模型",
        "    - YOLO/ONNX 布局识别模型",
        "    - XGBoost 段落连接模型"
    ])
    print("添加第19页：RAGFlowPdfParser 解析流程")

    # 第20页：MinerUParser 解析流程
    slide20 = add_content_slide(prs, "MinerUParser 解析流程（远程服务）", [
        "**1. __images__()** - 本地渲染PDF为图片（用于裁剪）",
        "**2. _run_mineru()** - HTTP POST /file_parse",
        "    - 发送PDF + 参数（backend, lang, method等）",
        "    - 下载ZIP结果并解压",
        "**3. _read_output()** - 查找 content_list.json",
        "    - 可选：查找 middle.json 增强表格位置",
        "**4. _transfer_to_sections()** - 转换为标准格式",
        "**5. _enhance_images_with_vlm()** - VLM图片描述增强（可选）",
        "",
        "**支持的 Backend：**",
        "    - pipeline: 传统多模型流水线（默认）",
        "    - vlm-transformers: 使用 HuggingFace Transformers",
        "    - vlm-vllm-engine: 本地 vLLM 引擎",
        "    - vlm-http-client: HTTP客户端远程VLM"
    ])
    print("添加第20页：MinerUParser 解析流程")

    # 第21页：OpenDataLoaderParser 解析流程
    slide21 = add_content_slide(prs, "OpenDataLoaderParser 解析流程（远程服务）", [
        "**1. __images__()** - 本地渲染PDF为图片（用于裁剪和坐标转换）",
        "**2. HTTP POST /v1/convert/file**",
        "    - 发送PDF + 参数（to_formats=[json, md]）",
        "**3. 解析响应** - 尝试多种嵌套层级提取 DoclingDocument",
        "**4. _normalize_docling_response()** - 标准化为内部格式",
        "    - 坐标转换：TOPLEFT -> bottom-left",
        "**5. _transfer_from_json()** - 转换为标准格式",
        "",
        "**Docling 格式特点：**",
        "    - prov 字段包含页面号和边界框",
        "    - table_cells 包含表格单元格信息",
        "    - 支持 HTML 重建和 Markdown 输出"
    ])
    print("添加第21页：OpenDataLoaderParser 解析流程")

    # 第22页：章节标题 - 定制点
    add_title_slide(prs, "DeepDoc 定制点", "如何扩展和定制解析器")
    print("添加第22页：章节标题 - 定制点")

    # 第23页：定制点1 - 添加新的解析器
    slide23 = add_content_slide(prs, "定制点 1：添加新的解析器", [
        "**步骤 1：创建解析器类（继承 RAGFlowPdfParser）**",
        "    class MyParser(RAGFlowPdfParser):",
        "        def __init__(self, api_url=\"\"):",
        "            self.api_url = api_url",
        "",
        "        def parse_pdf(self, filepath, binary=None, callback=None, **kwargs):",
        "            # 1. 调用外部服务或本地处理",
        "            # 2. 转换为标准格式：sections, tables",
        "            return sections, tables",
        "",
        "**步骤 2：在 naive.py 中注册**",
        "    def by_myparser(filename, binary=None, **kwargs):",
        "        pdf_parser = MyParser()",
        "        sections, tables = pdf_parser.parse_pdf(...)",
        "        return sections, tables, pdf_parser",
        "",
        "    PARSERS[\"myparser\"] = by_myparser",
        "",
        "**步骤 3：在配置中使用**",
        "    parser_config = {\"layout_recognize\": \"MyParser\"}"
    ])
    print("添加第23页：定制点1 - 添加新的解析器")

    # 第24页：定制点2 - 修改现有解析器行为
    slide24 = add_content_slide(prs, "定制点 2：修改现有解析器行为", [
        "**继承 RAGFlowPdfParser 并重写方法：**",
        "",
        "    class CustomPdfParser(RAGFlowPdfParser):",
        "        def __init__(self, **kwargs):",
        "            super().__init__(**kwargs)",
        "            # 自定义初始化",
        "",
        "        def __call__(self, fnm, need_image=True, zoomin=3, ...):",
        "            # 自定义解析流程",
        "            self.outlines = extract_pdf_outlines(fnm)",
        "            self.__images__(fnm, zoomin)",
        "            self._layouts_rec(zoomin)",
        "            ",
        "            # 自定义：添加额外的处理步骤",
        "            self._custom_processing()",
        "            ",
        "            self._table_transformer_job(zoomin)",
        "            self._text_merge()",
        "            ...",
        "",
        "**可重写的关键方法：**",
        "    __images__(), _layouts_rec(), _text_merge(),",
        "    _concat_downward(), _filter_forpages()"
    ])
    print("添加第24页：定制点2 - 修改现有解析器行为")

    # 第25页：定制点3 - 修改分块策略
    slide25 = add_content_slide(prs, "定制点 3：修改分块策略", [
        "**在 naive.py 的 chunk() 函数中定制：**",
        "",
        "    def chunk(filename, binary=None, ..., **kwargs):",
        "        # 获取解析结果",
        "        sections, tables, pdf_parser = parser(...)",
        "        ",
        "        # 自定义分块逻辑",
        "        custom_chunks = my_custom_chunking(sections, chunk_size=1024)",
        "        ",
        "        return custom_chunks",
        "",
        "**可调用的分块函数：**",
        "    - naive_merge() - 基础分块",
        "    - naive_merge_with_images() - 带图片的分块",
        "    - naive_merge_docx() - DOCX专用分块",
        "    - tokenize_chunks() - Tokenize分块",
        "    - doc_tokenize_chunks_with_images() - 带图片的Tokenize分块",
        "",
        "**配置参数：**",
        "    chunk_token_num: 分块大小（默认512）",
        "    delimiter: 分隔符（默认 \"\\n!?。；！？\"）"
    ])
    print("添加第25页：定制点3 - 修改分块策略")

    # 第26页：定制点4 - 配置解析器参数
    slide26 = add_content_slide(prs, "定制点 4：配置解析器参数", [
        "**通过 parser_config 传递参数：**",
        "",
        "    parser_config = {",
        "        \"layout_recognize\": \"DeepDOC\",  # 解析器选择",
        "        \"chunk_token_num\": 512,          # 分块大小",
        "        \"delimiter\": \"\\n!?。；！？\",    # 分隔符",
        "        \"table_context_size\": 0,          # 表格上下文",
        "        \"image_context_size\": 0,         # 图片上下文",
        "        \"pages\": [[1, 1000000]],          # 页面范围",
        "        \"analyze_hyperlink\": True,        # 分析超链接",
        "    }",
        "",
        "**MinerU 专用参数：**",
        "    mineru_lang: 语言（English/Chinese/...）",
        "    mineru_parse_method: 解析方法（auto/txt/ocr）",
        "    mineru_formula_enable: 公式识别",
        "    mineru_table_enable: 表格识别",
        "",
        "**OpenDataLoader 专用参数：**",
        "    hybrid: 混合处理模式",
        "    image_output: 图片输出方式",
        "    sanitize: 是否清理敏感内容"
    ])
    print("添加第26页：定制点4 - 配置解析器参数")

    # 第27页：定制点5 - 混合策略
    slide27 = add_content_slide(prs, "定制点 5：混合策略", [
        "**根据文档特征选择不同解析器：**",
        "",
        "    def by_smart(filename, binary=None, **kwargs):",
        "        if is_scan_pdf(binary):",
        "            # 扫描件 -> 使用 MinerU VLM",
        "            return by_mineru(filename, binary, **kwargs)",
        "        elif is_simple_pdf(binary):",
        "            # 简单PDF -> 使用本地 DeepDoc",
        "            return by_deepdoc(filename, binary, **kwargs)",
        "        else:",
        "            # 复杂文档 -> 使用 OpenDataLoader",
        "            return by_opendataloader(filename, binary, **kwargs)",
        "",
        "**实现思路：**",
        "    1. 检测PDF特征（扫描件/文本型/混合型）",
        "    2. 根据特征选择最优解析器",
        "    3. 可配置优先级和回退策略",
        "",
        "**优势：**",
        "    - 兼顾准确率和性能",
        "    - 降低单一解析器的局限性"
    ])
    print("添加第27页：定制点5 - 混合策略")

    # 第28页：关键文件位置
    slide28 = add_content_slide(prs, "关键文件位置", [
        "**入口层：**",
        "    rag/app/naive.py - 主入口，解析器调度",
        "",
        "**本地解析器：**",
        "    deepdoc/parser/pdf_parser.py - RAGFlowPdfParser",
        "    deepdoc/parser/vision/ocr.py - OCR引擎",
        "    deepdoc/parser/vision/layout_recognizer.py - 布局识别",
        "    deepdoc/parser/vision/table_structure_recognizer.py - 表格识别",
        "",
        "**远程解析器：**",
        "    deepdoc/parser/mineru_parser.py - MinerUParser",
        "    deepdoc/parser/opendataloader_parser.py - OpenDataLoaderParser",
        "    deepdoc/parser/docling_parser.py - DoclingParser",
        "",
        "**其他解析器：**",
        "    deepdoc/parser/tcadp_parser.py - TCADPParser（腾讯云）",
        "    deepdoc/parser/paddleocr_parser.py - PaddleOCRParser",
        "    deepdoc/parser/somark_parser.py - SoMarkParser",
        "",
        "**配置：**",
        "    api/db/db_models.py - 数据库模型（parser_id, parser_config）"
    ])
    print("添加第28页：关键文件位置")

    # 第29页：总结
    slide29 = add_content_slide(prs, "总结", [
        "**DeepDoc 解析链路设计灵活：**",
        "",
        "1. **入口层**：naive.py 的 chunk() 根据配置选择解析器",
        "2. **调度层**：PARSERS 字典将名称映射到具体函数",
        "3. **实现层**：每个解析器实现 parse_pdf() 返回标准格式",
        "4. **后处理层**：分块、tokenize、生成最终 chunks",
        "",
        "**核心定制点：**",
        "    - 添加新的解析器（继承 RAGFlowPdfParser）",
        "    - 修改现有解析器的行为（重写方法）",
        "    - 自定义分块策略",
        "    - 添加后处理步骤",
        "    - 通过配置参数调整行为",
        "    - 实现混合策略",
        "",
        "**设计哲学：**",
        "    统一的接口（sections + tables）",
        "    多样的实现（本地/远程/VLM）",
        "    灵活的配置（parser_config）"
    ])
    print("添加第29页：总结")

    # 保存PPT
    output_path = 'E:/桌面/work/Presentation/DeepDoc_Architecture.pptx'
    prs.save(output_path)
    print(f"\nPPT 已保存到: {output_path}")
    print(f"总页数: {len(prs.slides)} 页（原有15页 + 新增14页）")


if __name__ == "__main__":
    main()
