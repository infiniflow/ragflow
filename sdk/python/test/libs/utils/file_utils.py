#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#

import json

from docx import Document  # pip install python-docx
from openpyxl import Workbook  # pip install openpyxl
from PIL import Image, ImageDraw  # pip install Pillow
from pptx import Presentation  # pip install python-pptx
from reportlab.pdfgen import canvas  # pip install reportlab


def create_docx_file(path):
    doc = Document()
    doc.add_paragraph("这是一个测试 DOCX 文件。")
    doc.save(path)
    return path


def create_excel_file(path):
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "测试 Excel 文件"
    wb.save(path)
    return path


def create_ppt_file(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "测试 PPT 文件"
    prs.save(path)
    return path


def create_image_file(path):
    img = Image.new("RGB", (100, 100), color="blue")
    draw = ImageDraw.Draw(img)
    draw.text((10, 40), "Test", fill="white")
    img.save(path)
    return path


def create_pdf_file(path):
    if not isinstance(path, str):
        path = str(path)
    c = canvas.Canvas(path)
    c.drawString(100, 750, "测试 PDF 文件")
    c.save()
    return path


def create_txt_file(path):
    with open(path, "w", encoding="utf-8") as f:
        f.write("这是测试 TXT 文件的内容。")
    return path


def create_md_file(path):
    md_content = "# 测试 MD 文件\n\n这是一份 Markdown 格式的测试文件。"
    with open(path, "w", encoding="utf-8") as f:
        f.write(md_content)
    return path


def create_json_file(path):
    data = {"message": "这是测试 JSON 文件", "value": 123}
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    return path


def create_eml_file(path):
    eml_content = (
        "From: sender@example.com\n"
        "To: receiver@example.com\n"
        "Subject: 测试 EML 文件\n\n"
        "这是一封测试邮件的内容。\n"
    )
    with open(path, "w", encoding="utf-8") as f:
        f.write(eml_content)
    return path


def create_html_file(path):
    html_content = (
        "<html>\n"
        "<head><title>测试 HTML 文件</title></head>\n"
        "<body><h1>这是一个测试 HTML 文件</h1></body>\n"
        "</html>"
    )
    with open(path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return path
