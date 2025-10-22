#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#

import logging
from io import BytesIO
from pptx import Presentation
import tempfile
import subprocess
import os

class RAGFlowPptParser:
    def __init__(self):
        super().__init__()

    def __get_bulleted_text(self, paragraph):
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text

    def __extract(self, shape):
        try:
            # First try to get text content
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)

            # Safely get shape_type
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content
                if hasattr(shape, 'text'):
                    return shape.text.strip()
                return ""

            # Handle table
            if shape_type == 19:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append("; ".join([tb.cell(
                        0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
                return "\n".join(rows)

            # Handle group shape
            if shape_type == 6:
                texts = []
                for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def __call__(self, fnm, from_page, to_page, callback=None):
        os.environ[
            "LD_LIBRARY_PATH"] = "/usr/lib/libreoffice/program:" + os.environ.get(
            "LD_LIBRARY_PATH", "")
        if (isinstance(fnm, str) and fnm.lower().endswith('.ppt')) or (isinstance(fnm, bytes) and fnm[:8] == b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'):
            with tempfile.NamedTemporaryFile(suffix='.ppt') as tmp_ppt:
                if isinstance(fnm, bytes):
                    tmp_ppt.write(fnm)
                    tmp_ppt.flush()
                    ppt_path = tmp_ppt.name
                else:
                    ppt_path = fnm
                with tempfile.TemporaryDirectory() as tmp_dir:
                    cmd = [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "pptx",
                        "--outdir", tmp_dir,
                        ppt_path
                    ]
                    subprocess.run(cmd, check=True, capture_output=True)
                    pptx_name = os.path.splitext(os.path.basename(tmp_ppt.name))[0] + '.pptx'
                    pptx_path = os.path.join(tmp_dir, pptx_name)
                    ppt = Presentation(pptx_path)
        else:
            ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(
                BytesIO(fnm))

        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in sorted(
                    slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)):
                try:
                    txt = self.__extract(shape)
                    if txt:
                        texts.append(txt)
                except Exception as e:
                    logging.exception(e)
            txts.append("\n".join(texts))

        return txts
